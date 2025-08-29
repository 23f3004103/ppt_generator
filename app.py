from flask import Flask, request, send_file, render_template, jsonify
from werkzeug.utils import secure_filename
import os
import io
import tempfile
from pptx import Presentation
from pptx.util import Pt
import re
import requests

app = Flask(__name__, static_folder='static', template_folder='templates')
app.config['MAX_CONTENT_LENGTH'] = 30 * 1024 * 1024  # 30MB upload limit
ALLOWED_EXT = {'.pptx', '.potx'}

# Helpers

def allowed_file(filename):
    return os.path.splitext(filename)[1].lower() in ALLOWED_EXT


def split_text_to_sections(text, max_chars=1500):
    # naive split by paragraphs, chunk into reasonable sizes
    paras = [p.strip() for p in re.split(r"\n{2,}", text) if p.strip()]
    chunks = []
    cur = []
    cur_len = 0
    for p in paras:
        if cur_len + len(p) > max_chars and cur:
            chunks.append('\n\n'.join(cur))
            cur = [p]
            cur_len = len(p)
        else:
            cur.append(p)
            cur_len += len(p)
    if cur:
        chunks.append('\n\n'.join(cur))
    return chunks


def extract_template_styles(pptx_path):
    prs = Presentation(pptx_path)
    # Extract first slide's shapes as a simple style reference
    styles = {
        'layouts': [],
        'images': []
    }
    # grab some example images from the template (first slide)
    # Avoid slicing prs.slides (prs.slides[:1]) because python-pptx Slides.__getitem__
    # doesn't handle slice objects well in some versions and can return unexpected types.
    if len(prs.slides) > 0:
        slide = prs.slides[0]
        for shape in slide.shapes:
            try:
                if hasattr(shape, 'image') and shape.image:
                    img = shape.image
                    styles['images'].append(img.blob)
            except Exception:
                continue
    return styles


def create_presentation_from_chunks(chunks, template_path=None, guidance=None, refined=None):
    # If template provided, start from it; else create blank
    if template_path:
        prs = Presentation(template_path)
        template_styles = extract_template_styles(template_path)
    else:
        prs = Presentation()
        template_styles = {'images': []}
    # Simple mapping: one chunk -> 1-2 slides depending on length
    # If a template was provided, prefer to populate its existing slides in order
    existing_slides = list(prs.slides) if template_path else []

    # Helper functions for placing title/body into slides. Defined once to avoid nested
    # indentation and to keep behavior consistent across slides.
    def _set_title(sld, text):
        try:
            title_shape = sld.shapes.title
            if title_shape and getattr(title_shape, 'text_frame', None):
                tf = title_shape.text_frame
                try:
                    tf.text = text
                except Exception:
                    tf.clear()
                    tf.paragraphs[0].text = text
                return True
        except Exception:
            pass
        try:
            for ph in sld.placeholders:
                try:
                    if getattr(ph, 'is_title', False):
                        tf = ph.text_frame
                        try:
                            tf.text = text
                        except Exception:
                            tf.clear()
                            tf.paragraphs[0].text = text
                        return True
                except Exception:
                    continue
        except Exception:
            pass
        try:
            for shape in sld.shapes:
                try:
                    tf = getattr(shape, 'text_frame', None)
                    if tf is not None:
                        try:
                            tf.text = text
                        except Exception:
                            tf.clear()
                            tf.paragraphs[0].text = text
                        return True
                except Exception:
                    continue
        except Exception:
            pass
        return False

    def _set_body(sld, body_text):
        # prefer non-title placeholders
        try:
            for ph in sld.placeholders:
                try:
                    if not getattr(ph, 'is_title', False):
                        tf = ph.text_frame
                        try:
                            tf.text = '\n'.join(body_text.split('\n')[:20])
                        except Exception:
                            tf.clear()
                            for line in body_text.split('\n')[:20]:
                                p = tf.add_paragraph()
                                p.text = line
                                p.level = 0
                        return True
                except Exception:
                    continue
        except Exception:
            pass
        try:
            for shape in sld.shapes:
                try:
                    tf = getattr(shape, 'text_frame', None)
                    if tf is not None:
                        try:
                            tf.text = '\n'.join(body_text.split('\n')[:20])
                        except Exception:
                            tf.clear()
                            for line in body_text.split('\n')[:20]:
                                p = tf.add_paragraph()
                                p.text = line
                                p.level = 0
                        return True
                except Exception:
                    continue
        except Exception:
            pass
        try:
            from pptx.util import Inches
            left = top = Inches(1)
            width = Inches(8)
            height = Inches(4.5)
            txBox = sld.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = '\n'.join(body_text.split('\n')[:20])
            return True
        except Exception:
            return False

    def _set_body_with_bullets(sld, bullets_list):
        try:
            for ph in sld.placeholders:
                try:
                    if not getattr(ph, 'is_title', False):
                        tf = ph.text_frame
                        try:
                            tf.clear()
                        except Exception:
                            pass
                        for b in bullets_list[:10]:
                            p = tf.add_paragraph()
                            p.text = b
                            p.level = 0
                        return True
                except Exception:
                    continue
        except Exception:
            pass
        try:
            for shape in sld.shapes:
                try:
                    tf = getattr(shape, 'text_frame', None)
                    if tf is not None:
                        try:
                            tf.clear()
                        except Exception:
                            pass
                        for b in bullets_list[:10]:
                            p = tf.add_paragraph()
                            p.text = b
                            p.level = 0
                        return True
                except Exception:
                    continue
        except Exception:
            pass
        return False

    def _body_writer(sld, bullets_list, body_text):
        if bullets_list:
            return _set_body_with_bullets(sld, bullets_list)
        return _set_body(sld, body_text)

    for i, chunk in enumerate(chunks):
        # If refined data available, use that for title and bullets
        title = None
        bullets = None
        if refined and i < len(refined) and refined[i]:
            title = refined[i].get('title')
            bullets = refined[i].get('bullets')
            # create body from bullets (prefer bullets list)
            body = '\n'.join(bullets) if bullets else chunk
            if not title:
                # fallback: first short line
                lines = [l for l in re.split(r"\n+", chunk) if l.strip()]
                title = lines[0] if lines and len(lines[0]) < 100 else f"Slide {i+1}"
        else:
            lines = [l for l in re.split(r"\n+", chunk) if l.strip()]
            title = lines[0] if lines and len(lines[0]) < 100 else f"Slide {i+1}"
            body = '\n'.join(lines[1:]) if len(lines) > 1 else ''
        # If template has an existing slide at this index, populate it instead of adding
        if template_path and i < len(existing_slides):
            slide = existing_slides[i]
            # Clear existing textual content on the template slide so placeholder/sample
            # text from the template doesn't remain alongside our new content.
            for shape in list(slide.shapes):
                try:
                    # skip images
                    if hasattr(shape, 'image') and shape.image:
                        continue
                    if getattr(shape, 'has_text_frame', False):
                        tf = getattr(shape, 'text_frame', None)
                        if tf is not None:
                            try:
                                tf.clear()
                            except Exception:
                                try:
                                    tf.text = ''
                                except Exception:
                                    pass
                except Exception:
                    continue
        else:
            # choose a reasonable slide layout from the template if present
            try:
                # prefer a title layout for first slide
                if template_path:
                    if i == 0:
                        # find a layout with at least 1 placeholder
                        layout = next((l for l in prs.slide_layouts if len(list(l.placeholders)) >= 1), prs.slide_layouts[0])
                    else:
                        # find a layout with title+body (>=2 placeholders)
                        layout = next((l for l in prs.slide_layouts if len(list(l.placeholders)) >= 2), prs.slide_layouts[0])
                else:
                    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            except Exception:
                layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
        # Put title into the slide title placeholder if available, with fallbacks
        def _set_title(sld, text):
            try:
                title_shape = sld.shapes.title
                if title_shape and getattr(title_shape, 'text_frame', None):
                    tf = title_shape.text_frame
                    # avoid clearing to preserve template formatting where possible
                    try:
                        tf.text = text
                    except Exception:
                        tf.clear()
                        tf.paragraphs[0].text = text
                    return True
            except Exception:
                pass
            # fallback: look for a placeholder marked as title
            try:
                for ph in sld.placeholders:
                    try:
                        if getattr(ph, 'is_title', False):
                            tf = ph.text_frame
                            try:
                                tf.text = text
                            except Exception:
                                tf.clear()
                                tf.paragraphs[0].text = text
                            return True
                    except Exception:
                        continue
            except Exception:
                pass
            # final fallback: write into the first text-containing shape
            try:
                for shape in sld.shapes:
                    try:
                        tf = getattr(shape, 'text_frame', None)
                        if tf is not None:
                            tf.clear()
                            p = tf.paragraphs[0]
                            p.text = text
                            p.font.size = Pt(28)
                            return True
                    except Exception:
                        continue
            except Exception:
                pass
            return False
        
    print(f"Slide[{i}] - template_slide={bool(template_path and i < len(existing_slides))} title={repr(title)} bullets_count={len(bullets) if bullets else 0} body_len={len(body)}")
    title_ok = _set_title(slide, title or f"Slide {i+1}")
    print(f"  title_ok={title_ok}")

    # Put body into a main content placeholder if present, otherwise create textbox
    def _set_body(sld, body_text):
            # prefer non-title placeholders
            try:
                for ph in sld.placeholders:
                    try:
                        if not getattr(ph, 'is_title', False):
                            tf = ph.text_frame
                            # try to set text in one go to preserve placeholder styles
                            try:
                                tf.text = '\n'.join(body_text.split('\n')[:20])
                            except Exception:
                                tf.clear()
                                for line in body_text.split('\n')[:20]:
                                    p = tf.add_paragraph()
                                    p.text = line
                                    p.level = 0
                            return True
                    except Exception:
                        continue
            except Exception:
                pass
            # fallback: find any non-title text shape
            try:
                for shape in sld.shapes:
                    try:
                        tf = getattr(shape, 'text_frame', None)
                        if tf is not None:
                            try:
                                tf.text = '\n'.join(body_text.split('\n')[:20])
                            except Exception:
                                tf.clear()
                                for line in body_text.split('\n')[:20]:
                                    p = tf.add_paragraph()
                                    p.text = line
                                    p.level = 0
                            return True
                    except Exception:
                        continue
            except Exception:
                pass
            # final fallback: create a text box
            try:
                from pptx.util import Inches
                left = top = Inches(1)
                width = Inches(8)
                height = Inches(4.5)
                txBox = sld.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = '\n'.join(body_text.split('\n')[:20])
                return True
            except Exception:
                print('Warning: failed to place body text on slide', i)
                return False

    # prefer writing bullets as separate paragraphs when available
    def _set_body_with_bullets(sld, bullets_list):
        try:
            for ph in sld.placeholders:
                try:
                    if not getattr(ph, 'is_title', False):
                        tf = ph.text_frame
                        try:
                            tf.clear()
                        except Exception:
                            pass
                        for b in bullets_list[:10]:
                            p = tf.add_paragraph()
                            p.text = b
                            p.level = 0
                        return True
                except Exception:
                    continue
        except Exception:
            pass
        # fallback to any text shape
        try:
            for shape in sld.shapes:
                try:
                    tf = getattr(shape, 'text_frame', None)
                    if tf is not None:
                        try:
                            tf.clear()
                        except Exception:
                            pass
                        for b in bullets_list[:10]:
                            p = tf.add_paragraph()
                            p.text = b
                            p.level = 0
                        return True
                except Exception:
                    continue
        except Exception:
            pass
        return False

    def _body_writer(sld, bullets_list, body_text):
        if bullets_list:
            return _set_body_with_bullets(sld, bullets_list)
        return _set_body(sld, body_text)

    for i, chunk in enumerate(chunks):
        # If refined data available, use that for title and bullets
        title = None
        bullets = None
        if refined and i < len(refined) and refined[i]:
            title = refined[i].get('title')
            bullets = refined[i].get('bullets')
            # create body from bullets (prefer bullets list)
            body = '\n'.join(bullets) if bullets else chunk
            if not title:
                # fallback: first short line
                lines = [l for l in re.split(r"\n+", chunk) if l.strip()]
                title = lines[0] if lines and len(lines[0]) < 100 else f"Slide {i+1}"
        else:
            lines = [l for l in re.split(r"\n+", chunk) if l.strip()]
            title = lines[0] if lines and len(lines[0]) < 100 else f"Slide {i+1}"
            body = '\n'.join(lines[1:]) if len(lines) > 1 else ''

        # If template has an existing slide at this index, populate it instead of adding
        if template_path and i < len(existing_slides):
            slide = existing_slides[i]
        else:
            # choose a reasonable slide layout from the template if present
            try:
                # prefer a title layout for first slide
                if template_path:
                    if i == 0:
                        # find a layout with at least 1 placeholder
                        layout = next((l for l in prs.slide_layouts if len(list(l.placeholders)) >= 1), prs.slide_layouts[0])
                    else:
                        # find a layout with title+body (>=2 placeholders)
                        layout = next((l for l in prs.slide_layouts if len(list(l.placeholders)) >= 2), prs.slide_layouts[0])
                else:
                    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            except Exception:
                layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)

        print(f"Slide[{i}] - template_slide={bool(template_path and i < len(existing_slides))} title={repr(title)} bullets_count={len(bullets) if bullets else 0} body_len={len(body)}")
        title_ok = _set_title(slide, title or f"Slide {i+1}")
        print(f"  title_ok={title_ok}")

        body_ok = _body_writer(slide, bullets, body)
        print(f"  body_ok={body_ok}")

        # If we added a new slide, optionally reuse a template image as decorative asset.
        # For existing template slides we assume images are already present and skip re-adding.
        try:
            if (not (template_path and i < len(existing_slides))) and template_styles.get('images'):
                img_blob = template_styles['images'][0]
                img_stream = io.BytesIO(img_blob)
                from pptx.util import Inches
                pic = slide.shapes.add_picture(img_stream, Inches(6.5), Inches(4.5), width=Inches(2))
        except Exception:
            pass
    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio


def refine_chunk_with_llm(api_key, chunk, guidance='', model='gpt-3.5-turbo', base_url=None):
    """Call an OpenAI-compatible chat API to produce a concise title and bullet points for the chunk.
    Returns a dict: {'title':..., 'bullets':[...]} or None on failure."""
    if not api_key:
        print('LLM: no api_key provided, skipping')
        return None
    # lightweight log (do NOT log the api_key itself)
    print(f"LLM: calling refine_chunk_with_llm model={model} base_url={'provided' if base_url else 'openai'} chunk_len={len(chunk)}")
    prompt = (
        "You are a slide-writing assistant.\n"
        "Given the following content, produce a concise slide title (<=8 words) and 3-6 bullet points summarizing the most important items.\n"
        "Respond in JSON exactly as: {\"title\": \"...\", \"bullets\": [\"...\", ...]}\n\n"
        f"Guidance: {guidance}\n\nContent:\n{chunk}"
    )
    # Default to OpenAI-style chat completions unless base_url provided
    headers = {'Authorization': f'Bearer {api_key}', 'Content-Type': 'application/json'}
    payload = {
        'model': model,
        'messages': [
            {'role': 'system', 'content': 'You convert text into slide titles and bullets.'},
            {'role': 'user', 'content': prompt}
        ],
        'temperature': 0.2,
        'max_tokens': 300
    }
    # If base_url suggests Gemini generateContent, call that shape instead
    if base_url and 'geminiv1beta' in base_url:
        # build gemini payload
        gemini_payload = {'contents': [{'parts': [{'text': prompt}]}]}
        try:
            r = requests.post(base_url, json=gemini_payload, headers=headers, timeout=20)
            if r.status_code != 200:
                print('LLM: gemini returned', r.status_code, r.text[:200])
                return None
            j = r.json()
            content = j.get('candidates', [{}])[0].get('content', {}).get('parts', [''])[0]
            print('LLM: gemini response length', len(content))
        except Exception as e:
            print('LLM: gemini request failed:', str(e))
            return None
    else:
        # default OpenAI-style
        r = None
        try:
            url = base_url.rstrip('/') + '/chat/completions' if base_url else 'https://api.openai.com/v1/chat/completions'
            r = requests.post(url, json=payload, headers=headers, timeout=15)
            if r.status_code != 200:
                print('LLM: chat completions returned', r.status_code, (r.text or '')[:200])
                return None
            j = r.json()
            content = j.get('choices', [{}])[0].get('message', {}).get('content', '')
            print('LLM: chat/completions response length', len(content))
        except Exception as e:
            print('LLM: chat/completions request failed:', str(e))
            return None
    # if a custom base_url is passed via guidance variable (we overload guidance tuple), skip - but default uses OpenAI
    # We'll accept guidance as a dict in the future; for now use the default OpenAI endpoint.
    url = f"https://api.openai.com/v1/chat/completions"
    # try to parse JSON out of content
    try:
        import json as _json
        # content might contain markdown or text around JSON; extract first JSON object
        m = re.search(r"\{[\s\S]*\}", content)
        if m:
            parsed = _json.loads(m.group(0))
            title = parsed.get('title', '').strip()
            bullets = parsed.get('bullets', []) or []
            bullets = [b.strip() for b in bullets][:8]
            print('LLM: parsed JSON title="%s" bullets=%d' % (title, len(bullets)))
            return {'title': title or None, 'bullets': bullets}
    except Exception:
        return None
    except Exception:
        return None
    return None


@app.route('/llm_test', methods=['POST'])
def llm_test():
    """Test endpoint: returns the LLM-refined title/bullets for a single piece of text.
    POST form fields: api_key, text, guidance (optional), base_url (optional), model (optional)
    """
    api_key = request.form.get('api_key', '')
    text = request.form.get('text', '')
    guidance = request.form.get('guidance', '')
    base_url = request.form.get('base_url') or None
    model = request.form.get('model') or 'gpt-3.5-turbo'
    if not api_key:
        return jsonify({'error': 'no api_key provided'}), 400
    if not text.strip():
        return jsonify({'error': 'no text provided'}), 400
    res = refine_chunk_with_llm(api_key, text, guidance=guidance, model=model, base_url=base_url)
    if res is None:
        return jsonify({'error': 'LLM call failed or returned unparsable output'}), 500
    return jsonify({'result': res})


@app.route('/')
def index():
    return render_template('index.html')


@app.route('/generate', methods=['POST'])
def generate():
    # Do not log API keys or user text
    text = request.form.get('text', '')
    guidance = request.form.get('guidance', '')
    api_key = request.form.get('api_key', '')
    template = request.files.get('template')
    base_url = request.form.get('base_url') or None
    model = request.form.get('model') or 'gpt-3.5-turbo'

    if not text.strip():
        return jsonify({'error': 'No input text provided'}), 400

    template_path = None
    if template and template.filename:
        if not allowed_file(template.filename):
            return jsonify({'error': 'Invalid template file type'}), 400
        fname = secure_filename(template.filename)
        tmpdir = tempfile.gettempdir()
        # ensure tmpdir exists (should on most systems)
        os.makedirs(tmpdir, exist_ok=True)
        save_path = os.path.join(tmpdir, fname)
        template.save(save_path)
        template_path = save_path

    # Split text into chunks
    chunks = split_text_to_sections(text)
    # Cap to a reasonable number of slides
    MAX_SLIDES = 12
    if len(chunks) > MAX_SLIDES:
        chunks = chunks[:MAX_SLIDES]

    # If API key provided, refine each chunk with LLM to get title + bullets
    refined = None
    if api_key:
        refined = []
        for c in chunks:
            res = refine_chunk_with_llm(api_key, c, guidance=guidance, model=model, base_url=base_url)
            refined.append(res)

    # TODO: call LLM to refine titles/bullets. For now, naive mapping.
    prs_bio = create_presentation_from_chunks(chunks, template_path=template_path, guidance=guidance, refined=refined)

    return send_file(prs_bio, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation', as_attachment=True, download_name='generated_presentation.pptx')


if __name__ == '__main__':
    app.run(debug=True)
